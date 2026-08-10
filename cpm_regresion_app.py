# ─────────────────────────────────────────────────────────────────────────────
# Common Progress Measures — Regression Analysis Dashboard
# Built by Sian Reilly, Data & Intelligence Analyst
# Strategy & Intelligence, Westminster City Council
#
# pip install streamlit plotly pandas scikit-learn openpyxl python-pptx kaleido
# Run: streamlit run cpm_regression_app.py
# ─────────────────────────────────────────────────────────────────────────────

import streamlit as st
import plotly.express as px
import plotly.graph_objects as go
import pandas as pd
import numpy as np
import io, os, warnings
from sklearn.linear_model import LassoCV, Lasso
from sklearn.preprocessing import StandardScaler
from sklearn.model_selection import LeaveOneOut, cross_val_predict
from sklearn.metrics import r2_score, mean_absolute_error
from pptx import Presentation
from pptx.util import Inches, Pt

warnings.filterwarnings("ignore")

# ─── COLOURS & SETUP ──────────────────────────────────────────────────────────
PILLAR_COLOURS = {
    "Money, local economy, jobs and training": "#6D2E46",
    "Crime and safety": "#A26769",
    "Education, community and connection": "#C89B7B",
    "Health, wellbeing and healthcare": "#4A7C59",
    "Housing and homelessness": "#2C5F7C",
    "Neighbourhood and environment": "#8B6F47",
}
BERRY, DUSTY, ACCENT, INK, MUTED, GREEN = "#6D2E46", "#A26769", "#C89B7B", "#2E1018", "#7A5C62", "#4A7C59"
YEAR_SHEETS = {"CPM_2024_2025_Scores": "2024-25", "CPM_2023_2024_Scores": "2023-24"}
WARD_NAME_MAP = {"Knightsbridge and Belgravia": "Knightsbridge & Belgravia"}
LE_F, LE_M, LE_O = "Life expectancy, female", "Life expectancy, male", "Life expectancy, overall"
TARGETS = {"male": LE_M, "female": LE_F, "overall": LE_O}
TARGET_DISPLAY = {"male": "👨 Male", "female": "👩 Female", "overall": "👥 Overall (avg of both)"}

# ─── HELPERS ──────────────────────────────────────────────────────────────────
def shorten(name, maxlen=45):
    s = (name.replace("Proportion of ", "").replace("Percentage of ", "% ")
         .replace("% of residents who ", "").replace("% of residents with ", "")
         .replace("Number of ", "No. ").replace("Average ", "Avg ")
         .replace("Placement of ", "").replace("- medium (4-5) or high (6-9)", ""))
    return s[:maxlen-3] + "..." if len(s) > maxlen else s

def pptx_btn(fig, cid):
    try: img = fig.to_image(format="png", width=1200, height=700, scale=2)
    except: return
    prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.shapes.add_picture(io.BytesIO(img), Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.0))
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    st.download_button("⬇ Download as slide", buf, f"{cid}.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation", key=f"dl_{cid}")

def style(fig):
    fig.update_layout(font_family="Arial", title_font_size=16, title_font_color="#333",
        plot_bgcolor="white", paper_bgcolor="white", margin=dict(l=40, r=20, t=60, b=40),
        legend=dict(orientation="h", yanchor="bottom", y=1.02))
    fig.update_xaxes(showgrid=False, linecolor="#ccc")
    fig.update_yaxes(gridcolor="#eee", linecolor="white")
    return fig

def chart(fig, cid, cap=""):
    style(fig); st.plotly_chart(fig, use_container_width=True, key=cid)
    if cap: st.caption(cap)
    pptx_btn(fig, cid)

# ─── DATA ─────────────────────────────────────────────────────────────────────
_PILLAR = {}

@st.cache_data
def parse_sheet(filepath, sheet_name):
    raw = pd.read_excel(filepath, sheet_name=sheet_name, header=None)
    pillars, metrics, pols, ctypes = raw.iloc[1].ffill(), raw.iloc[2].ffill(), raw.iloc[3].ffill(), raw.iloc[5]
    ws = raw.iloc[6:].copy(); ws = ws[ws.iloc[:,0].notna()].reset_index(drop=True)
    recs = []
    for j in range(1, raw.shape[1]):
        ct = str(ctypes.iloc[j]).strip() if pd.notna(ctypes.iloc[j]) else ""
        if "min" not in ct.lower(): continue
        m = str(metrics.iloc[j]).strip() if pd.notna(metrics.iloc[j]) else ""
        p = str(pillars.iloc[j]).strip() if pd.notna(pillars.iloc[j]) else ""
        pol = str(pols.iloc[j]).strip() if pd.notna(pols.iloc[j]) else ""
        if not m or m == "Metric:": continue
        outlier = bool("outlier" in m.lower() or "removed" in m.lower())
        for i in range(len(ws)):
            ward = WARD_NAME_MAP.get(str(ws.iloc[i,0]).strip(), str(ws.iloc[i,0]).strip())
            v = ws.iloc[i, j]
            if pd.notna(v):
                try: recs.append({"Ward": ward, "Metric": m, "Pillar": p, "Polarity": pol,
                    "Score": float(v), "Outlier_Version": outlier})
                except: pass
    return pd.DataFrame(recs)

@st.cache_data
def load_years(fp):
    d = {}
    for sh, lb in YEAR_SHEETS.items():
        try: df = parse_sheet(fp, sh); df["Year"] = lb; d[lb] = df
        except: pass
    return d

@st.cache_data
def load_imd(fp): df = pd.read_excel(fp, sheet_name="IMD_2025_Ward"); df["WD24NM"] = df["WD24NM"].replace(WARD_NAME_MAP); return df

@st.cache_data
def load_overview(fp): return pd.read_excel(fp, sheet_name="Dataset_Overview")

def make_wide(df, no_outliers=True):
    if no_outliers: df = df[~df["Outlier_Version"]]
    w = df.pivot_table(index="Ward", columns="Metric", values="Score", aggfunc="first")
    if "Westminster" in w.index: w = w.drop("Westminster")
    # Flexibly find LE columns by searching for partial matches
    le_male_col = _find_le_col(w.columns, "male")
    le_female_col = _find_le_col(w.columns, "female")
    # Rename to standard names so the rest of the app works
    rename_map = {}
    if le_male_col and le_male_col != LE_M: rename_map[le_male_col] = LE_M
    if le_female_col and le_female_col != LE_F: rename_map[le_female_col] = LE_F
    if rename_map: w = w.rename(columns=rename_map)
    # Create overall as average of male + female
    if LE_M in w.columns and LE_F in w.columns: w[LE_O] = w[[LE_M, LE_F]].mean(axis=1)
    return w

def _find_le_col(columns, sex):
    """Flexibly find life expectancy column by partial match."""
    sex_l = sex.lower()
    for c in columns:
        cl = c.lower().strip()
        if "life" in cl and "expectancy" in cl:
            # "female" contains "male" so we must check carefully
            if sex_l == "male" and "female" in cl:
                continue  # skip female columns when looking for male
            if sex_l in cl:
                return c
    return None

def run_lasso(wide, target):
    le = [c for c in wide.columns if "life expectancy" in c.lower()]
    feat = [c for c in wide.columns if c not in le]
    if target not in wide.columns: return None
    mdf = wide[[target]+feat].dropna(subset=[target])
    mdf = mdf.dropna(axis=1, thresh=int(len(mdf)*0.6)).fillna(mdf.median())
    y, X = mdf[target], mdf[[c for c in feat if c in mdf.columns]]
    if len(X.columns)==0: return None
    # Enforce numeric and check for zero variance
    try:
        X = X.astype(np.float64)
        y = y.astype(np.float64)
        if float(y.std()) == 0: return None
    except (ValueError, TypeError):
        return None
    try:
        Xs = StandardScaler().fit_transform(X)
        loo = LeaveOneOut()
        lcv = LassoCV(cv=loo, random_state=42, max_iter=10000, n_alphas=100); lcv.fit(Xs, y)
        mdl = Lasso(alpha=lcv.alpha_, max_iter=10000); mdl.fit(Xs, y)
        yloo = cross_val_predict(Lasso(alpha=lcv.alpha_, max_iter=10000), Xs, y, cv=loo)
    except Exception:
        return None
    coefs = pd.DataFrame({"Metric": X.columns, "Short": [shorten(c) for c in X.columns],
        "Coefficient": mdl.coef_, "Abs": np.abs(mdl.coef_),
        "Pillar": [_PILLAR.get(c,"Unknown") for c in X.columns]}).sort_values("Abs", ascending=False)
    preds = pd.DataFrame({"Ward": mdf.index, "Actual": y.values, "Predicted": yloo,
        "Residual": y.values - yloo}).sort_values("Residual")
    return {"coefs": coefs, "preds": preds, "r2": r2_score(y, yloo), "r2in": r2_score(y, mdl.predict(Xs)),
        "mae": mean_absolute_error(y, yloo), "nsel": int(np.sum(mdl.coef_!=0)),
        "ntot": len(X.columns), "X": X, "y": y, "target": target}

def run_all(wide):
    out = {}
    for k, v in TARGETS.items():
        if v in wide.columns:
            r = run_lasso(wide, v)
            if r:
                out[k] = r
    return out

# ─── PAGE ─────────────────────────────────────────────────────────────────────
st.set_page_config(page_title="CPM Regression Analysis", page_icon="📊", layout="wide", initial_sidebar_state="expanded")

with st.sidebar:
    st.header("⚙️ Settings")
    uploaded = st.file_uploader("Upload normalised_outputs.xlsx", type=["xlsx"])
    fp = uploaded
    if fp is None:
        for f in ["normalised_outputs.xlsx", "Dummy_Datasets.xlsx", "Core_Datasets.xlsx"]:
            if os.path.exists(f): fp = f; break
    st.markdown("---")
    no_outliers = not st.checkbox("Use outlier-removed versions", value=False,
        help="Default OFF — keeps all 18 wards. Turning this on uses cleaner data but may drop some wards.")
    st.markdown("---")
    st.caption("Built by Sian Reilly · Data & Intelligence · WCC")

if fp is None:
    st.title("Common Progress Measures — Regression Analysis")
    st.warning("Upload your Excel file using the sidebar to get started.")
    st.stop()

try:
    yd = load_years(fp); imd = load_imd(fp); overview = load_overview(fp)
except Exception as e:
    st.error(f"Error loading data: {e}"); st.stop()

for lb, df in yd.items():
    for _, r in df.drop_duplicates("Metric").iterrows(): _PILLAR[r["Metric"]] = r["Pillar"]

yrs = sorted(yd.keys()); latest = yrs[-1]
wides, results = {}, {}
for lb in yrs:
    w = make_wide(yd[lb], no_outliers); wides[lb] = w; results[lb] = run_all(w)

# Diagnostic: show what was found in sidebar
with st.sidebar:
    st.markdown("---")
    st.markdown("**🔍 Data diagnostic**")
    for lb in yrs:
        w = wides[lb]
        le_found = [c for c in w.columns if "life expectancy" in c.lower()]
        n_results = len(results.get(lb, {}))
        st.caption(f"{lb}: {len(w)} wards, {len(w.columns)} cols, LE cols: {le_found}, models: {n_results}")
    if not any(results.values()):
        all_metrics = []
        for lb, df in yd.items():
            all_metrics.extend(df["Metric"].unique().tolist())
        le_like = [m for m in set(all_metrics) if "life" in m.lower() or "le " in m.lower() or m.lower().startswith("le")]
        if le_like:
            st.warning(f"Found LE-like metrics in raw data: {le_like}")
        else:
            st.warning(f"No LE metrics found. Sample metrics: {list(set(all_metrics))[:5]}")

# ─── MAIN ─────────────────────────────────────────────────────────────────────
st.title("Common Progress Measures — Regression Analysis")
st.markdown("*Built by Sian Reilly, Data & Intelligence Analyst, Strategy & Intelligence, Westminster City Council*")
st.caption("This dashboard tests how well our CPM indicators explain life expectancy across Westminster's 18 wards, "
           "identifies which indicators matter most, and flags wards where outcomes are unexpectedly good or bad.")

tabs = st.tabs(["🏠 How to Read This", "🔬 Data Overview (EDA)", "📊 What Explains LE?",
    "🔍 Unexpected Wards", "📅 Year-on-Year", "🗺️ Deprivation & LE", "📈 Explore the Data"])

# ═══════════════════════════════════════════════════════════════════════════════
# TAB 0 — HOW TO READ THIS
# ═══════════════════════════════════════════════════════════════════════════════
with tabs[0]:
    st.header("What this dashboard does — and how to read it")
    st.markdown("""
    This tool applies two analytical techniques to the Common Progress Measures to answer
    questions that the existing ward scoring can't:

    ### What is LASSO regression?

    Imagine you have 29 CPM indicators and you want to know which ones actually matter for
    life expectancy. You could look at each one individually, but the problem is that many
    indicators move together — wards with high poverty tend to also have high worklessness,
    poor housing, and worse health outcomes. So which indicator is actually doing the explaining,
    and which ones are just along for the ride?

    **LASSO regression** solves this. It's a technique that:
    1. Looks at all 29 indicators at once
    2. Works out which combination best predicts life expectancy
    3. **Automatically drops** indicators that don't add anything beyond what the others already capture
    4. Gives each remaining indicator a **coefficient** — a number showing how strongly it relates to LE

    The key advantage over looking at simple correlations is that LASSO handles **overlapping
    indicators** (the technical term is multicollinearity). If child poverty and workless households
    both predict LE but are highly correlated with each other, LASSO will typically keep one and
    drop the other rather than double-counting.

    The output is a short list of indicators that genuinely matter, with everything else filtered out.

    ### What is residual analysis?

    The current CPM scoring tells you "Church Street scores low on anxiety." But it doesn't tell
    you whether that's *surprising*. A ward that's deprived across the board will naturally score
    low on lots of things — that's expected, not news.

    **Residual analysis** asks a different question: given everything else we know about a ward,
    is this particular score better or worse than we'd *expect*?

    The **residual** is simply: **Actual value minus Predicted value**.

    - A **large negative residual** means life expectancy is worse than the ward's overall CPM
      profile would suggest — something specific may need attention
    - A **large positive residual** means LE is better than expected — something is going right

    This shifts the conversation from "which wards are worst" (which we already know) to
    "where are the specific, unexpected gaps that targeted action could address."

    ### What is "Overall LE"?

    The dashboard runs the analysis for **male LE, female LE, and overall LE** (which is the
    average of the male and female normalised scores). This isn't a separate data source — it's
    a combined view that smooths out sex-specific patterns.

    ### What does R² mean?

    R² (R-squared) is a number between 0 and 1 that answers: **what share of the variation
    in life expectancy can our CPM indicators explain?**

    - **R² = 0.75** → the indicators explain 75% of why LE differs across wards. Strong.
    - **R² = 0.30** → they only explain 30%. We're missing something.

    The R² shown here uses **Leave-One-Out cross-validation** — each ward is left out in turn
    and predicted from the remaining 17. This gives an honest estimate rather than an
    overly optimistic one.

    ### Important caveats

    - **18 wards is a small sample.** Results are directional signals, not precise measurements.
      Adding or removing a single ward can change results.
    - **Life expectancy is lagging.** It's measured over 5-year periods, so it may not reflect
      recent changes in CPM indicators.
    - **Air quality has reverse polarity** (Damian's point) — pollution is worst in the wealthiest
      areas (West End, Marylebone). The normalisation handles the scoring direction, but the
      underlying confound remains.
    - **Correlation isn't causation.** These indicators *co-occur* with LE differences; they don't
      necessarily *cause* them.
    """)

    st.info("💡 Every chart has a **⬇ Download as slide** button so you can drop it straight into a deck.")

# ═══════════════════════════════════════════════════════════════════════════════
# TAB 1 — EXPLORATORY DATA ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
with tabs[1]:
    st.header("Data Overview — Exploratory Data Analysis")
    st.markdown("""
    Before running any models, it's important to understand what we're working with.
    This section profiles the dataset following the standard EDA pipeline:
    **Profile → Identify issues → Visualise → Interpret.**
    """)

    w = wides[latest]
    le_cols = [c for c in w.columns if "life expectancy" in c.lower()]
    feat_cols = [c for c in w.columns if c not in le_cols]

    # ── Profile ──
    st.subheader("1. Dataset profile")
    c1, c2, c3 = st.columns(3)
    c1.metric("Wards", len(w))
    c2.metric("CPM indicators", len(feat_cols))
    c3.metric("LE measures", f"{len(le_cols)} ({', '.join([c.split(', ')[1] for c in le_cols if ', ' in c])})")

    st.markdown(f"**Year:** {latest} · **Geography:** Ward level · **Values:** Min-max normalised scores (0–1 scale, higher = better)")

    # Pillar breakdown
    pillar_counts = pd.Series([_PILLAR.get(c, "Unknown") for c in feat_cols]).value_counts()
    fig_pillars = px.bar(x=pillar_counts.values, y=pillar_counts.index, orientation="h",
        color=pillar_counts.index, color_discrete_map=PILLAR_COLOURS,
        labels={"x": "Number of indicators", "y": ""})
    fig_pillars.update_layout(title="Indicators by CPM pillar", showlegend=False, height=350)
    chart(fig_pillars, "eda_pillars", "The six CPM pillars, showing how many indicators each contributes to the model.")

    # ── Missingness ──
    st.subheader("2. Missing data")
    missing = w[feat_cols].isna().sum()
    missing_pct = (missing / len(w) * 100).round(1)
    missing_df = pd.DataFrame({"Indicator": [shorten(c) for c in missing.index],
        "Missing count": missing.values, "Missing %": missing_pct.values})
    missing_any = missing_df[missing_df["Missing count"] > 0]

    if len(missing_any) == 0:
        st.success("No missing values across any indicator — the dataset is complete.")
    else:
        st.warning(f"{len(missing_any)} indicator(s) have missing data:")
        st.dataframe(missing_any.sort_values("Missing %", ascending=False), use_container_width=True, hide_index=True)
        st.caption("Where data is missing, the model fills gaps with the median ward score for that indicator. "
                   "With 18 wards, this is preferable to dropping wards or indicators entirely.")

    # ── Distribution ──
    st.subheader("3. How are the scores distributed?")
    st.markdown("Since all indicators are on a 0–1 normalised scale, we can compare their distributions directly. "
                "Indicators where all wards cluster together (low spread) will be less useful for the regression "
                "than those with clear variation across wards.")

    spreads = w[feat_cols].std().sort_values(ascending=False)
    spread_df = pd.DataFrame({"Indicator": [shorten(c) for c in spreads.index],
        "Std Dev": spreads.values, "Pillar": [_PILLAR.get(c,"") for c in spreads.index]})
    fig_spread = go.Figure(go.Bar(y=spread_df["Indicator"], x=spread_df["Std Dev"], orientation="h",
        marker_color=[PILLAR_COLOURS.get(p, MUTED) for p in spread_df["Pillar"]]))
    fig_spread.update_layout(title="Variation across wards: which indicators spread most?",
        xaxis_title="Standard deviation (higher = more variation between wards)",
        height=max(500, len(spread_df)*20+100))
    chart(fig_spread, "eda_spread",
        "Indicators with very low spread (short bars) may not help distinguish between wards.")

    # ── Correlation with LE ──
    st.subheader("4. Initial correlations with life expectancy")
    st.markdown("Before running the regression, a simple correlation check shows which indicators "
                "have the strongest *individual* relationship with LE. But remember — individual "
                "correlations can be misleading when indicators overlap. That's why we use LASSO.")

    eda_target = st.selectbox("Show correlations for", list(TARGET_DISPLAY.keys()),
        format_func=lambda k: TARGET_DISPLAY[k], key="eda_corr_target")
    tgt_col = TARGETS[eda_target]
    if tgt_col in w.columns:
        corrs = w[feat_cols].corrwith(w[tgt_col]).sort_values()
        corr_df = pd.DataFrame({"Indicator": [shorten(c) for c in corrs.index],
            "Correlation": corrs.values, "Pillar": [_PILLAR.get(c,"") for c in corrs.index]})
        fig_corr = go.Figure(go.Bar(y=corr_df["Indicator"], x=corr_df["Correlation"], orientation="h",
            marker_color=[PILLAR_COLOURS.get(p, MUTED) for p in corr_df["Pillar"]]))
        fig_corr.update_layout(title=f"Individual correlation of each indicator with {TARGET_DISPLAY[eda_target]} LE",
            xaxis_title="Pearson correlation (−1 to +1)", height=max(500, len(corr_df)*20+100))
        chart(fig_corr, f"eda_corr_{eda_target}",
            "Positive = higher score → higher LE. Negative = higher score → lower LE. "
            "But these don't account for overlap between indicators — LASSO does.")

    # ── Pillar correlation heatmap ──
    st.subheader("5. Do the CPM pillars overlap?")
    st.markdown("If two pillars are highly correlated, their indicators are largely measuring the same "
                "thing. LASSO handles this automatically, so you don't need to manually group or remove "
                "indicators — but it's useful to see where the redundancy sits.")

    pillar_avgs = pd.DataFrame()
    for p in sorted(set(_PILLAR.values())):
        pcols = [c for c in feat_cols if _PILLAR.get(c) == p]
        if pcols: pillar_avgs[p.split(",")[0] if "," in p else p[:25]] = w[pcols].mean(axis=1)
    if len(pillar_avgs.columns) > 1:
        fig_heat = px.imshow(pillar_avgs.corr(), text_auto=".2f",
            color_continuous_scale=["white", DUSTY, BERRY], zmin=-1, zmax=1, aspect="auto")
        fig_heat.update_layout(title="Correlation between CPM pillars (ward-level averages)", height=450)
        chart(fig_heat, "eda_pillar_heat",
            "Values close to 1.0 mean two pillars move together — their indicators overlap. "
            "LASSO deals with this by keeping one representative and dropping the rest.")

# ═══════════════════════════════════════════════════════════════════════════════
# TAB 2 — WHAT EXPLAINS LE (LASSO)
# ═══════════════════════════════════════════════════════════════════════════════
with tabs[2]:
    st.header("What explains life expectancy across Westminster's wards?")
    st.markdown("The LASSO regression runs for **all three** measures simultaneously: male, female, "
                "and overall (average of both). This lets us see which indicators consistently matter "
                "and which are specific to one sex.")

    # Summary cards
    cols = st.columns(3)
    for i, (k, disp) in enumerate(TARGET_DISPLAY.items()):
        with cols[i]:
            r = results[latest].get(k)
            if r:
                st.subheader(disp)
                st.metric("R² (cross-validated)", f"{r['r2']:.2f}",
                    help="What share of LE variation the model explains. Higher = better.")
                st.metric("Indicators kept", f"{r['nsel']} out of {r['ntot']}",
                    help="How many indicators LASSO selected vs how many it started with.")
                st.metric("Average error", f"{r['mae']:.3f}",
                    help="Typical prediction error in normalised score units.")
            else:
                st.subheader(disp); st.warning("No data")

    st.markdown("---")

    # Detailed view
    detail_k = st.selectbox("Explore in detail", list(TARGET_DISPLAY.keys()),
        format_func=lambda k: TARGET_DISPLAY[k], key="lasso_detail")
    r = results[latest].get(detail_k)

    if r:
        sel = r["coefs"][r["coefs"]["Abs"] > 0].sort_values("Coefficient")
        dropped = r["coefs"][r["coefs"]["Abs"] == 0]

        if len(sel) > 0:
            st.subheader(f"Indicators selected by LASSO for {TARGET_DISPLAY[detail_k]} LE ({latest})")
            st.markdown("Each bar shows the **coefficient** — how strongly that indicator relates to LE "
                        "after accounting for everything else. Longer bar = stronger relationship.")
            fig = go.Figure()
            for _, row in sel.iterrows():
                fig.add_trace(go.Bar(y=[row["Short"]], x=[row["Coefficient"]], orientation="h",
                    marker_color=PILLAR_COLOURS.get(row["Pillar"], MUTED), showlegend=False,
                    hovertemplate=f"<b>{row['Short']}</b><br>Pillar: {row['Pillar']}<br>Coefficient: {row['Coefficient']:.3f}<extra></extra>"))
            fig.update_layout(title=f"LASSO kept {len(sel)} indicators — these are the ones that matter",
                xaxis_title="Coefficient (positive = associated with higher LE)", height=max(350, len(sel)*40+100))
            chart(fig, f"lasso_{detail_k}",
                "Positive bars: wards scoring higher on this indicator tend to have higher LE. "
                "Negative bars: the opposite. Coloured by CPM pillar.")
            legend = " · ".join([f'<span style="color:{PILLAR_COLOURS.get(p, MUTED)}">■</span> {p}' for p in sel["Pillar"].unique()])
            st.markdown(f"**Pillar key:** {legend}", unsafe_allow_html=True)
        else:
            st.warning("LASSO couldn't find reliable predictors at this sample size.")

        with st.expander(f"📋 {len(dropped)} indicators LASSO dropped (click to see)"):
            st.markdown("These were **zeroed out** — they don't add explanatory value beyond "
                        "what the selected indicators already capture. This doesn't mean they're "
                        "unimportant — just that other indicators already cover the same ground.")
            st.dataframe(dropped[["Short", "Pillar"]].rename(columns={"Short": "Indicator"}),
                use_container_width=True, hide_index=True)

        # Actual vs predicted
        st.subheader("How well does the model predict each ward?")
        st.markdown("Each dot is a ward. If the model were perfect, every dot would sit on the "
                    "dashed line. Dots above the line have *higher* LE than predicted; below = lower.")
        pdf = r["preds"]
        fig2 = px.scatter(pdf, x="Predicted", y="Actual", text="Ward", color_discrete_sequence=[BERRY])
        mn, mx = min(pdf["Actual"].min(), pdf["Predicted"].min())-0.05, max(pdf["Actual"].max(), pdf["Predicted"].max())+0.05
        fig2.add_trace(go.Scatter(x=[mn,mx], y=[mn,mx], mode="lines", line=dict(dash="dash", color="#ccc"), name="Perfect prediction"))
        fig2.update_traces(textposition="top center", selector=dict(mode="markers+text"))
        fig2.update_layout(title=f"Actual vs Predicted — {TARGET_DISPLAY[detail_k]} LE",
            xaxis_title="What the model predicts", yaxis_title="What actually happened")
        chart(fig2, f"avp_{detail_k}")

    # Cross-target comparison
    st.markdown("---")
    st.subheader("Which indicators appear across male, female, and overall?")
    st.markdown("Indicators selected for **all three** are the strongest signals. "
                "Those appearing in only one may be noise.")
    cross = []
    for k in ["male", "female", "overall"]:
        r = results[latest].get(k)
        if r:
            for _, row in r["coefs"][r["coefs"]["Abs"]>0].iterrows():
                cross.append({"Indicator": row["Short"], "Target": TARGET_DISPLAY[k], "Coefficient": row["Coefficient"]})
    if cross:
        cdf = pd.DataFrame(cross)
        pivot = cdf.pivot_table(index="Indicator", columns="Target", values="Coefficient", aggfunc="first").fillna(0)
        cnt = (pivot != 0).sum(axis=1).sort_values(ascending=False)
        all3 = cnt[cnt==3].index.tolist()
        two = cnt[cnt==2].index.tolist()
        if all3: st.success(f"**Selected for all three:** {', '.join(all3)}")
        if two: st.info(f"**Selected for two:** {', '.join(two)}")

        fig_c = px.bar(cdf, x="Coefficient", y="Indicator", color="Target", orientation="h",
            barmode="group", color_discrete_sequence=[BERRY, DUSTY, ACCENT])
        fig_c.update_layout(title="How each indicator relates to male, female, and overall LE",
            height=max(400, len(cdf["Indicator"].unique())*30+100))
        chart(fig_c, "cross_target")

# ═══════════════════════════════════════════════════════════════════════════════
# TAB 3 — RESIDUAL ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
with tabs[3]:
    st.header("Which wards have unexpectedly high or low life expectancy?")
    st.markdown("""
    The residual is the gap between what the model **predicts** for a ward (given its CPM scores)
    and what **actually** happens. A large gap means something unusual is going on.

    - **Red bars (negative residual):** LE is *worse* than the ward's CPM profile would suggest.
      These wards may benefit from targeted investigation or intervention.
    - **Green bars (positive residual):** LE is *better* than expected. Something is going right
      that the CPM indicators don't fully capture.
    """)

    rc1, rc2 = st.columns(2)
    res_k = rc1.selectbox("LE measure", list(TARGET_DISPLAY.keys()), format_func=lambda k: TARGET_DISPLAY[k], key="res_k")
    res_yr = rc2.selectbox("Year", yrs, index=len(yrs)-1, key="res_yr")
    r = results.get(res_yr, {}).get(res_k)

    if r:
        pdf = r["preds"]
        fig3 = go.Figure(go.Bar(y=pdf["Ward"], x=pdf["Residual"], orientation="h",
            marker_color=[BERRY if v<0 else GREEN for v in pdf["Residual"]],
            hovertemplate="<b>%{y}</b><br>Residual: %{x:.3f}<extra></extra>"))
        fig3.update_layout(title=f"Where does {TARGET_DISPLAY[res_k]} LE diverge from expectations? ({res_yr})",
            xaxis_title="Residual (actual minus predicted)", height=max(450, len(pdf)*30+80))
        chart(fig3, f"resid_{res_k}_{res_yr}")

        detail = pdf.copy().round(3)
        detail["Interpretation"] = detail["Residual"].apply(
            lambda v: "⚠️ LE worse than expected" if v < -0.1 else ("✅ LE better than expected" if v > 0.1 else "➖ Roughly as expected"))
        st.dataframe(detail, use_container_width=True, hide_index=True)

        st.markdown("---")
        st.subheader("Drill down: what's unusual about a specific ward?")
        st.markdown("Select a ward to see which CPM indicators are most different from the Westminster "
                    "average. This helps explain *why* the ward's LE is unexpectedly high or low.")
        X = r["X"]
        wc = st.selectbox("Select ward", list(X.index), key=f"drill_{res_k}_{res_yr}")
        if wc in X.index:
            diff = (X.loc[wc] - X.mean()).sort_values()
            ddf = pd.DataFrame({"Indicator": [shorten(c) for c in diff.index], "vs_Avg": diff.values,
                "Ward_Score": X.loc[wc].values, "Westminster_Avg": X.mean().values,
                "Pillar": [_PILLAR.get(c,"") for c in diff.index]})
            top = pd.concat([ddf.head(5), ddf.tail(5)])
            fig4 = go.Figure()
            for _, row in top.iterrows():
                fig4.add_trace(go.Bar(y=[row["Indicator"]], x=[row["vs_Avg"]], orientation="h",
                    marker_color=PILLAR_COLOURS.get(row["Pillar"], MUTED), showlegend=False,
                    hovertemplate=f"<b>{row['Indicator']}</b><br>{wc}: {row['Ward_Score']:.2f}<br>Avg: {row['Westminster_Avg']:.2f}<extra></extra>"))
            fig4.update_layout(title=f"{wc}: indicators most different from the Westminster average",
                xaxis_title="Difference (negative = worse than average, positive = better)", height=420)
            chart(fig4, f"drill_{wc.replace(' ','_')}_{res_k}")
    else:
        st.warning("No results for this combination.")

# ═══════════════════════════════════════════════════════════════════════════════
# TAB 4 — YEAR-ON-YEAR
# ═══════════════════════════════════════════════════════════════════════════════
with tabs[4]:
    st.header("Year-on-Year: 2023-24 vs 2024-25")
    st.markdown("Running the same model on two years of data lets us check **stability** — "
                "indicators that appear in both years are more trustworthy than those that flip.")

    if len(yrs) < 2:
        st.warning("Need both years."); st.stop()

    y1, y2 = yrs[0], yrs[1]
    yoy_k = st.selectbox("LE measure", list(TARGET_DISPLAY.keys()), format_func=lambda k: TARGET_DISPLAY[k], key="yoy_k")
    r1, r2r = results.get(y1,{}).get(yoy_k), results.get(y2,{}).get(yoy_k)

    if r1 and r2r:
        mc1, mc2 = st.columns(2)
        for col, lb, r in [(mc1, y1, r1), (mc2, y2, r2r)]:
            col.subheader(lb); col.metric("R²", f"{r['r2']:.2f}"); col.metric("Selected", f"{r['nsel']}/{r['ntot']}")

        st.markdown("---")
        st.subheader("Feature stability: which indicators were selected in both years?")
        ca = []
        for lb, r in [(y1, r1), (y2, r2r)]:
            c = r["coefs"][["Short","Coefficient","Pillar"]].copy(); c["Year"] = lb; ca.append(c)
        ca = pd.concat(ca)
        pv = ca.pivot_table(index="Short", columns="Year", values="Coefficient", aggfunc="first").fillna(0)
        act = pv[(pv.abs()>0).any(axis=1)]
        if len(act) > 0:
            fig5 = go.Figure()
            for yr in [y1, y2]:
                if yr in act.columns:
                    fig5.add_trace(go.Bar(y=act.index, x=act[yr], name=yr, orientation="h",
                        marker_color=BERRY if yr==y1 else DUSTY))
            fig5.update_layout(title=f"Indicator coefficients: {y1} vs {y2} ({TARGET_DISPLAY[yoy_k]})",
                barmode="group", xaxis_title="Coefficient", height=max(400, len(act)*30+100))
            chart(fig5, f"yoy_coefs_{yoy_k}")

            both = pv[(pv.abs()>0).all(axis=1)].index.tolist()
            one = pv[(pv.abs()>0).any(axis=1) & ~(pv.abs()>0).all(axis=1)].index.tolist()
            if both: st.success(f"**Stable (both years):** {', '.join(both)}")
            if one: st.warning(f"**One year only (less reliable):** {', '.join(one)}")

        st.markdown("---")
        st.subheader("Did any wards shift position between years?")
        st.markdown("This shows how each ward's residual changed. A ward moving into the red "
                    "is doing *worse relative to expectations* than it was last year.")
        p1 = r1["preds"][["Ward","Residual"]].rename(columns={"Residual": y1})
        p2 = r2r["preds"][["Ward","Residual"]].rename(columns={"Residual": y2})
        rc = p1.merge(p2, on="Ward"); rc["Change"] = rc[y2] - rc[y1]; rc = rc.sort_values("Change")
        fig6 = go.Figure(go.Bar(y=rc["Ward"], x=rc["Change"], orientation="h",
            marker_color=[BERRY if c<0 else GREEN for c in rc["Change"]]))
        fig6.update_layout(title=f"Residual change {y1} → {y2}: {TARGET_DISPLAY[yoy_k]}",
            xaxis_title="Change (negative = declined vs expectations)", height=max(450, len(rc)*28+80))
        chart(fig6, f"resid_chg_{yoy_k}")

        st.markdown("---")
        st.subheader("Individual indicator movement")
        common = sorted(set(wides[y1].columns) & set(wides[y2].columns))
        common = [m for m in common if "life expectancy" not in m.lower()]
        if common:
            ch = st.selectbox("Select indicator", common, format_func=shorten, key="yoy_metric")
            sc = pd.DataFrame({y1: wides[y1].get(ch), y2: wides[y2].get(ch)}).dropna()
            sc["Change"] = sc[y2] - sc[y1]; sc = sc.sort_values("Change")
            fig7 = go.Figure(go.Bar(y=sc.index, x=sc["Change"], orientation="h",
                marker_color=[BERRY if c<0 else GREEN for c in sc["Change"]]))
            fig7.update_layout(title=f"{shorten(ch)}: ward-level change {y1} → {y2}",
                xaxis_title="Change in normalised score (positive = improved)", height=max(400, len(sc)*28+80))
            chart(fig7, "metric_move")
    else:
        st.warning("Results not available for both years.")

# ═══════════════════════════════════════════════════════════════════════════════
# TAB 5 — DEPRIVATION & LE
# ═══════════════════════════════════════════════════════════════════════════════
with tabs[5]:
    st.header("Deprivation and Life Expectancy")
    st.markdown("The Index of Multiple Deprivation (IMD 2025) measures how deprived each ward is across "
                "income, employment, education, health, crime, housing, and living environment. "
                "A **higher IMD score = more deprived**. Decile 1 = most deprived 10% nationally.")

    # IMD overview
    imd_s = imd.sort_values(" IMD25 IMD Score", ascending=False)
    fig_imd = go.Figure(go.Bar(y=imd_s["WD24NM"], x=imd_s[" IMD25 IMD Score"], orientation="h", marker_color=BERRY))
    fig_imd.update_layout(title="Church Street is Westminster's most deprived ward (IMD 2025)",
        xaxis_title="IMD Score (higher = more deprived)", height=520)
    chart(fig_imd, "imd_overall")

    st.markdown("---")

    # IMD vs LE scatter
    st.subheader("Does higher deprivation mean lower life expectancy?")
    st.markdown("We'd expect a negative relationship — more deprived wards should have lower LE. "
                "Wards far from the trendline are the 'unexpected' ones.")
    imd_le_k = st.selectbox("LE measure", list(TARGET_DISPLAY.keys()), format_func=lambda k: TARGET_DISPLAY[k], key="imd_le")
    le_col = TARGETS[imd_le_k]
    if le_col in wides[latest].columns:
        les = wides[latest][le_col].copy()
        les.name = "LE_Score"
        imds = imd.set_index("WD24NM")[" IMD25 IMD Score"].copy()
        imds.name = "IMD_Score"
        sdf = pd.DataFrame({"Ward": les.index, "LE_Score": les.values, "IMD_Score": imds.reindex(les.index).values}).dropna()
        fig_sc = px.scatter(sdf, x="IMD_Score", y="LE_Score", text="Ward",
            color_discrete_sequence=[BERRY], trendline="ols")
        fig_sc.update_traces(textposition="top center", selector=dict(mode="markers+text"))
        fig_sc.update_layout(title=f"IMD deprivation vs {TARGET_DISPLAY[imd_le_k]} LE ({latest})",
            xaxis_title="IMD Score (higher = more deprived)", yaxis_title="LE (normalised, higher = better)")
        chart(fig_sc, f"imd_le_{imd_le_k}")

    # LE change by deprivation
    if len(yrs) >= 2:
        st.markdown("---")
        st.subheader("Are deprived wards falling further behind?")
        st.markdown("This plots each ward's LE *change* between years against its deprivation level. "
                    "A downward trend would mean the gap is widening.")
        chg_k = st.selectbox("LE measure", list(TARGET_DISPLAY.keys()), format_func=lambda k: TARGET_DISPLAY[k], key="le_chg")
        lc = TARGETS[chg_k]
        if lc in wides[y1].columns and lc in wides[y2].columns:
            lchg = pd.DataFrame({y1: wides[y1][lc], y2: wides[y2][lc]}).dropna()
            lchg["LE_Change"] = lchg[y2] - lchg[y1]
            lchg["IMD"] = lchg.index.map(imd.set_index("WD24NM")[" IMD25 IMD Score"])
            lchg = lchg.dropna(subset=["IMD"]).sort_values("IMD", ascending=False)

            fig_chg = go.Figure(go.Bar(y=lchg.index, x=lchg["LE_Change"], orientation="h",
                marker=dict(color=lchg["IMD"], colorscale=[[0,GREEN],[1,BERRY]],
                    colorbar=dict(title="IMD")),
                hovertemplate="<b>%{y}</b><br>LE Change: %{x:.3f}<br>IMD: %{marker.color:.1f}<extra></extra>"))
            fig_chg.update_layout(title=f"LE change ({y1} → {y2}), sorted by deprivation (most deprived at top)",
                xaxis_title="Change in normalised LE score", height=max(450, len(lchg)*30+80))
            chart(fig_chg, f"le_chg_imd_{chg_k}",
                "Darker red bars = more deprived wards. If deprived wards are declining while affluent ones improve, "
                "the inequality gap is widening.")

            lchg_plot = lchg.reset_index()
            ward_col = lchg_plot.columns[0]  # whatever the index was named
            fig_sc2 = px.scatter(lchg_plot, x="IMD", y="LE_Change", text=ward_col,
                color_discrete_sequence=[BERRY], trendline="ols")
            fig_sc2.update_traces(textposition="top center", selector=dict(mode="markers+text"))
            fig_sc2.add_hline(y=0, line_dash="dash", line_color="#ccc")
            fig_sc2.update_layout(title=f"Is LE change linked to deprivation? ({TARGET_DISPLAY[chg_k]})",
                xaxis_title="IMD Score", yaxis_title=f"LE change {y1} → {y2}")
            chart(fig_sc2, f"imd_chg_scatter_{chg_k}",
                "Downward trendline = deprived wards declining faster. Flat = inequality stable.")

    st.markdown("---")

    # Domain comparison and decile table
    st.subheader("IMD domain profiles")
    doms = ["Income","Employment","Education, Skills and Training","Health","Crime","Barriers to Housing and Services","Living Environment"]
    dr = []
    for d in doms:
        c = f" IMD25 {d} Score"
        if c in imd.columns:
            for _, row in imd.iterrows(): dr.append({"Ward": row["WD24NM"], "Domain": d, "Score": row[c]})
    ddf = pd.DataFrame(dr)
    ws = st.multiselect("Compare wards", imd["WD24NM"].tolist(),
        default=["Church Street", "Westbourne", "Knightsbridge & Belgravia", "Regent's Park"])
    if ws:
        fig_d = px.bar(ddf[ddf["Ward"].isin(ws)], x="Score", y="Domain", color="Ward", orientation="h",
            barmode="group", color_discrete_sequence=px.colors.qualitative.Set2)
        fig_d.update_layout(title="IMD domain profiles", xaxis_title="Score (higher = more deprived)", height=420)
        chart(fig_d, "imd_doms")

    st.subheader("Decile overview")
    dc = [c for c in imd.columns if "Decile" in c]
    dd = imd[["WD24NM"]+dc].copy()
    dd.columns = ["Ward"] + [c.replace(" IMD25 ","").replace(" Decile","") for c in dc]
    dd = dd.sort_values("IMD")
    try:
        st.dataframe(dd.style.background_gradient(cmap="RdYlGn_r", subset=dd.columns[1:]),
            use_container_width=True, hide_index=True)
    except Exception:
        st.dataframe(dd, use_container_width=True, hide_index=True)
    st.caption("Decile 1 (red) = most deprived 10% nationally. Decile 10 (green) = least deprived.")

# ═══════════════════════════════════════════════════════════════════════════════
# TAB 6 — EXPLORE
# ═══════════════════════════════════════════════════════════════════════════════
with tabs[6]:
    st.header("Explore the Data")
    st.markdown("Use this tab to explore individual relationships between any indicator and life expectancy.")

    exp_k = st.selectbox("LE measure", list(TARGET_DISPLAY.keys()), format_func=lambda k: TARGET_DISPLAY[k], key="exp_k")
    r = results[latest].get(exp_k)
    if r:
        sc_opts = sorted(r["X"].columns)
        sc_ch = st.selectbox("Select indicator", sc_opts, format_func=shorten, key="exp_scatter")
        edf = pd.DataFrame({"Indicator": r["X"][sc_ch], "LE": r["y"], "Ward": r["X"].index})
        fig10 = px.scatter(edf, x="Indicator", y="LE", text="Ward", color_discrete_sequence=[BERRY], trendline="ols")
        fig10.update_traces(textposition="top center", selector=dict(mode="markers+text"))
        fig10.update_layout(title=f"{shorten(sc_ch)} vs {TARGET_DISPLAY[exp_k]} LE",
            xaxis_title=shorten(sc_ch), yaxis_title="LE (normalised)")
        chart(fig10, f"explore_{exp_k}")

    st.markdown("---")
    st.subheader("Full dataset")
    st.dataframe(overview, use_container_width=True, hide_index=True)
    st.subheader(f"Ward × Indicator matrix ({latest})")
    st.dataframe(wides[latest].round(3), use_container_width=True)
